from pathlib import Path
from pptx import Presentation


downloads = Path.home() / "Downloads"
pptx_files = [p for p in downloads.glob("*.pptx") if "26.4.13" in p.name]
if not pptx_files:
    raise SystemExit("No matching pptx found")

p = pptx_files[0]
print(f"PPTX: {p}")
prs = Presentation(p)
print(f"SLIDES: {len(prs.slides)}")

start = 1
end = len(prs.slides)
if len(__import__("sys").argv) == 3:
    start = int(__import__("sys").argv[1])
    end = int(__import__("sys").argv[2])

for i, slide in enumerate(prs.slides, 1):
    if i < start or i > end:
        continue
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text = " ".join(line.strip() for line in shape.text.splitlines() if line.strip())
            if text:
                texts.append(text)
    print(f"\n--- SLIDE {i} ---")
    for text in texts:
        print(text)
